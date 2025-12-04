import io
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from fpdf import FPDF

def extract_text_from_file(uploaded_file):
    """Lê o arquivo baseado na extensão e retorna texto puro."""
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    try:
        if file_type in ['txt', 'md']:
            return uploaded_file.read().decode("utf-8")
        
        elif file_type == 'pdf':
            reader = PdfReader(uploaded_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_type == 'docx':
            doc = Document(uploaded_file)
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif file_type == 'pptx':
            prs = Presentation(uploaded_file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif file_type in ['xlsx', 'xls']:
            df = pd.read_excel(uploaded_file)
            return df.to_string()
            
        else:
            return "Formato não suportado."
    except Exception as e:
        return f"Erro ao ler arquivo: {str(e)}"

def generate_docx(text):
    doc = Document()
    doc.add_heading('História de Usuário (CWS)', 0)
    for paragraph in text.split('\n'):
        doc.add_paragraph(paragraph)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()

def generate_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    # Tratamento básico para evitar erro de encoding no FPDF
    text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 10, text)
    return pdf.output(dest='S').encode('latin-1')